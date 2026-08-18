"""Document parsers.

Each parser reports a confidence and the structural elements it could not
handle. The platform never claims full-fidelity extraction: a scanned PDF
without OCR configured returns low confidence and names the gap, so a
downstream answer can be qualified rather than silently wrong.
"""

from __future__ import annotations

import csv
import io
import json
import re
from dataclasses import dataclass, field
from typing import Any

from agentic_os.core.errors import NotImplementedCapability, ValidationError


@dataclass(slots=True)
class ParseResult:
    text: str
    confidence: float
    mime_type: str
    page_count: int | None = None
    unsupported_elements: list[str] = field(default_factory=list)
    structured: dict[str, Any] | None = None
    parser: str = ""

    def to_dict(self) -> dict[str, Any]:
        return {
            "confidence": self.confidence,
            "mime_type": self.mime_type,
            "page_count": self.page_count,
            "unsupported_elements": self.unsupported_elements,
            "parser": self.parser,
            "characters": len(self.text),
        }


SUPPORTED_MIME_TYPES = {
    "text/plain",
    "text/markdown",
    "text/csv",
    "text/html",
    "application/json",
    "application/xml",
    "text/xml",
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

#: Formats the platform recognises but does not extract. Declared so the gap is
#: visible in the ingestion record rather than surfacing as an empty document.
DECLARED_UNSUPPORTED = {
    "image/png": "OCR is not configured; register an OCR adapter to ingest images",
    "image/jpeg": "OCR is not configured; register an OCR adapter to ingest images",
    "image/tiff": "OCR is not configured; register an OCR adapter to ingest images",
    "audio/mpeg": "audio transcription is not configured",
    "audio/wav": "audio transcription is not configured",
    "video/mp4": "video transcription is not configured",
    "application/vnd.ms-outlook": "email export parsing is not implemented",
}


def parse_text(data: bytes, mime_type: str = "text/plain") -> ParseResult:
    text = data.decode("utf-8", errors="replace")
    replacements = text.count("�")
    confidence = 1.0 if replacements == 0 else max(0.5, 1.0 - replacements / max(1, len(text)))
    unsupported = ["undecodable byte sequences"] if replacements else []
    return ParseResult(
        text=text,
        confidence=round(confidence, 3),
        mime_type=mime_type,
        unsupported_elements=unsupported,
        parser="text",
    )


def parse_csv(data: bytes) -> ParseResult:
    raw = data.decode("utf-8", errors="replace")
    reader = csv.DictReader(io.StringIO(raw))
    rows = list(reader)
    headers = reader.fieldnames or []
    lines = [" | ".join(headers)]
    for row in rows[:5000]:
        lines.append(" | ".join(str(row.get(h, "") or "") for h in headers))
    return ParseResult(
        text="\n".join(lines),
        confidence=1.0 if headers else 0.4,
        mime_type="text/csv",
        unsupported_elements=[] if headers else ["no header row detected"],
        structured={"headers": headers, "rows": rows, "row_count": len(rows)},
        parser="csv",
    )


def parse_json(data: bytes) -> ParseResult:
    try:
        payload = json.loads(data.decode("utf-8"))
    except (json.JSONDecodeError, UnicodeDecodeError) as exc:
        raise ValidationError(f"invalid JSON document: {exc}") from exc
    return ParseResult(
        text=json.dumps(payload, indent=2, ensure_ascii=False),
        confidence=1.0,
        mime_type="application/json",
        structured={"payload": payload},
        parser="json",
    )


_TAG = re.compile(r"<[^>]+>")
_SCRIPT_STYLE = re.compile(r"<(script|style)\b[^>]*>.*?</\1>", re.IGNORECASE | re.DOTALL)


def parse_html(data: bytes) -> ParseResult:
    raw = data.decode("utf-8", errors="replace")
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(raw, "html.parser")
        for element in soup(["script", "style", "noscript"]):
            element.decompose()
        text = soup.get_text("\n")
        parser = "beautifulsoup"
        confidence = 0.95
    except ImportError:
        text = _TAG.sub(" ", _SCRIPT_STYLE.sub(" ", raw))
        parser = "regex_fallback"
        confidence = 0.7
    text = re.sub(r"\n{3,}", "\n\n", re.sub(r"[ \t]{2,}", " ", text)).strip()
    return ParseResult(
        text=text,
        confidence=confidence,
        mime_type="text/html",
        unsupported_elements=["embedded scripts and styles are discarded"],
        parser=parser,
    )


def parse_xml(data: bytes) -> ParseResult:
    # defusedxml is not a dependency, so parsing uses a text extraction that
    # never resolves entities — this cannot be used for an XXE attack.
    raw = data.decode("utf-8", errors="replace")
    if re.search(r"<!ENTITY", raw, re.IGNORECASE) or "<!DOCTYPE" in raw.upper():
        raise ValidationError(
            "XML documents containing a DOCTYPE or entity declaration are rejected "
            "(external entity expansion risk)"
        )
    text = _TAG.sub(" ", raw)
    text = re.sub(r"\s{2,}", " ", text).strip()
    return ParseResult(
        text=text,
        confidence=0.8,
        mime_type="application/xml",
        unsupported_elements=["attributes and namespaces are not preserved"],
        parser="xml_text",
    )


def parse_pdf(data: bytes) -> ParseResult:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - optional extra
        raise NotImplementedCapability(
            "PDF parsing requires the 'parsers' extra: pip install -e '.[parsers]'"
        ) from exc

    reader = PdfReader(io.BytesIO(data))
    pages: list[str] = []
    empty_pages = 0
    for page in reader.pages:
        extracted = (page.extract_text() or "").strip()
        if not extracted:
            empty_pages += 1
        pages.append(extracted)

    total = len(pages) or 1
    text_ratio = 1 - (empty_pages / total)
    unsupported: list[str] = []
    if empty_pages:
        unsupported.append(
            f"{empty_pages} of {total} pages yielded no extractable text "
            "(likely scanned images; OCR is not configured)"
        )
    return ParseResult(
        text="\n\n".join(f"[page {i + 1}]\n{p}" for i, p in enumerate(pages) if p),
        confidence=round(max(0.1, text_ratio * 0.95), 3),
        mime_type="application/pdf",
        page_count=total,
        unsupported_elements=unsupported,
        parser="pypdf",
    )


def parse_docx(data: bytes) -> ParseResult:
    try:
        import docx
    except ImportError as exc:  # pragma: no cover - optional extra
        raise NotImplementedCapability(
            "DOCX parsing requires the 'parsers' extra: pip install -e '.[parsers]'"
        ) from exc

    document = docx.Document(io.BytesIO(data))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    unsupported = []
    if document.inline_shapes:
        unsupported.append(f"{len(document.inline_shapes)} embedded images not extracted")
    return ParseResult(
        text="\n".join(parts),
        confidence=0.93,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        unsupported_elements=unsupported + ["headers, footers and footnotes are not extracted"],
        parser="python-docx",
    )


def parse_xlsx(data: bytes) -> ParseResult:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover - optional extra
        raise NotImplementedCapability(
            "XLSX parsing requires the 'parsers' extra: pip install -e '.[parsers]'"
        ) from exc

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines: list[str] = []
    sheets: dict[str, list[dict]] = {}
    for sheet in workbook.worksheets:
        lines.append(f"# Sheet: {sheet.title}")
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        headers = [str(h) if h is not None else f"col_{i}" for i, h in enumerate(rows[0])]
        records: list[dict] = []
        for row in rows[1:5001]:
            record = {headers[i]: row[i] for i in range(min(len(headers), len(row)))}
            records.append(record)
            lines.append(" | ".join("" if v is None else str(v) for v in row))
        sheets[sheet.title] = records
    return ParseResult(
        text="\n".join(lines),
        confidence=0.9,
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        unsupported_elements=["formulas are read as cached values; charts are not extracted"],
        structured={"sheets": sheets},
        parser="openpyxl",
    )


def parse_pptx(data: bytes) -> ParseResult:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - optional extra
        raise NotImplementedCapability(
            "PPTX parsing requires the 'parsers' extra: pip install -e '.[parsers]'"
        ) from exc

    presentation = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for number, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[slide {number}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return ParseResult(
        text="\n".join(p for p in parts if p.strip()),
        confidence=0.88,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        page_count=len(presentation.slides),
        unsupported_elements=["speaker notes, images and diagrams are not extracted"],
        parser="python-pptx",
    )


_DISPATCH = {
    "text/plain": lambda d: parse_text(d, "text/plain"),
    "text/markdown": lambda d: parse_text(d, "text/markdown"),
    "text/csv": parse_csv,
    "application/json": parse_json,
    "text/html": parse_html,
    "application/xml": parse_xml,
    "text/xml": parse_xml,
    "application/pdf": parse_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": parse_docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": parse_xlsx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": parse_pptx,
}


def parse(data: bytes, mime_type: str) -> ParseResult:
    """Parse a document. Raises NotImplementedCapability for declared gaps."""
    if mime_type in DECLARED_UNSUPPORTED:
        raise NotImplementedCapability(DECLARED_UNSUPPORTED[mime_type], details={"mime_type": mime_type})
    handler = _DISPATCH.get(mime_type)
    if handler is None:
        raise NotImplementedCapability(
            f"no parser registered for mime type '{mime_type}'",
            details={"supported": sorted(SUPPORTED_MIME_TYPES)},
        )
    return handler(data)
